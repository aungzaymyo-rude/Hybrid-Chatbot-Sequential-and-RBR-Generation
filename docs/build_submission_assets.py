# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
import re
from textwrap import wrap

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches as PptInches, Pt as PptPt
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / 'docs'
FIG = DOCS / 'figures'
PREVIEW_DIR = DOCS / 'presentation_previews'
PREVIEW_DIR.mkdir(exist_ok=True)

DISSERTATION_MD = DOCS / 'dissertation_supervisor_ready.md'
DISSERTATION_DOCX = DOCS / 'dissertation_supervisor_ready.docx'
CHECKLIST_MD = DOCS / 'submission_checklist.md'
CHECKLIST_DOCX = DOCS / 'submission_checklist.docx'
MANUAL_MD = DOCS / 'user_manual_guide.md'
MANUAL_DOCX = DOCS / 'user_manual_guide.docx'
PPTX_OUT = DOCS / 'dissertation_presentation.pptx'

FIGURE_MAP = {
    'Figure 3.1': FIG / 'figure_3_1_hybrid_architecture.png',
    'Figure 3.2': FIG / 'figure_3_2_mlops_lifecycle.png',
    'Figure 3.3': FIG / 'figure_3_3_chat_ui.png',
    'Figure 4.1': FIG / 'figure_4_1_metric_comparison.png',
    'Figure 4.2': FIG / 'figure_4_2_training_curves.png',
    'Figure 4.3': FIG / 'figure_4_3_dataset_and_splits.png',
    'Figure 4.4': FIG / 'figure_4_4_testing_summary.png',
    'Figure 4.5': FIG / 'figure_4_5_admin_overview.png',
    'Figure 4.6': FIG / 'figure_4_7_admin_trace.png',
}

ACCENT = RGBColor(13, 84, 78)
ACCENT_LIGHT = RGBColor(225, 241, 238)
TEXT_DARK = RGBColor(24, 39, 40)
TEXT_MUTED = RGBColor(81, 98, 100)
BG_LIGHT = RGBColor(247, 250, 249)
BG_DARK = RGBColor(14, 49, 47)


def clean_inline_markdown(text: str) -> str:
    text = text.replace('**', '')
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    return text.replace('`', '')


def set_default_font(style, name: str = 'Arial', size: int = 11) -> None:
    style.font.name = name
    style.font.size = Pt(size)
    style._element.rPr.rFonts.set(qn('w:ascii'), name)
    style._element.rPr.rFonts.set(qn('w:hAnsi'), name)


def set_document_layout(doc: Document) -> None:
    for section in doc.sections:
        section.left_margin = Cm(3)
        section.right_margin = Cm(2)
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
    set_default_font(doc.styles['Normal'])
    set_default_font(doc.styles['Title'], size=20)
    set_default_font(doc.styles['Heading 1'], size=16)
    set_default_font(doc.styles['Heading 2'], size=13)
    if 'Heading 3' in doc.styles:
        set_default_font(doc.styles['Heading 3'], size=11)


def set_para_spacing(paragraph, line_15: bool = True) -> None:
    fmt = paragraph.paragraph_format
    fmt.space_after = Pt(6)
    fmt.space_before = Pt(0)
    if line_15:
        fmt.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    else:
        fmt.line_spacing = 1.0


def add_page_field(paragraph) -> None:
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fld = OxmlElement('w:fldSimple')
    fld.set(qn('w:instr'), 'PAGE')
    run = OxmlElement('w:r')
    text = OxmlElement('w:t')
    text.text = '1'
    run.append(text)
    fld.append(run)
    paragraph._p.append(fld)


def restart_page_numbering(section) -> None:
    sect_pr = section._sectPr
    pg_num_type = sect_pr.find(qn('w:pgNumType'))
    if pg_num_type is None:
        pg_num_type = OxmlElement('w:pgNumType')
        sect_pr.append(pg_num_type)
    pg_num_type.set(qn('w:start'), '1')


def add_toc(paragraph) -> None:
    fld = OxmlElement('w:fldSimple')
    fld.set(qn('w:instr'), 'TOC \\o "1-3" \\h \\z \\u')
    run = OxmlElement('w:r')
    text = OxmlElement('w:t')
    text.text = 'Right-click and update field in Word to generate the table of contents.'
    run.append(text)
    fld.append(run)
    paragraph._p.append(fld)


def add_body_paragraph(doc: Document, text: str, style: str | None = None, single_space: bool = False):
    p = doc.add_paragraph(style=style)
    p.add_run(clean_inline_markdown(text))
    set_para_spacing(p, line_15=not single_space)
    return p


def add_caption(doc: Document, text: str):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(text)
    run.bold = True
    run.italic = True
    run.font.name = 'Arial'
    run.font.size = Pt(10)
    set_para_spacing(p, line_15=False)
    return p


def add_figure(doc: Document, image_path: Path, caption: str, width: Inches = Inches(6.3)):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(image_path), width=width)
    set_para_spacing(p, line_15=False)
    add_caption(doc, caption)


def add_metrics_table(doc: Document):
    add_caption(doc, 'Table 4.1 Final test-set metrics for the general and report model profiles')
    table = doc.add_table(rows=1, cols=6)
    table.style = 'Table Grid'
    headers = ['Model', 'Accuracy', 'Macro F1', 'Weighted F1', 'Macro Precision', 'Macro Recall']
    for cell, header in zip(table.rows[0].cells, headers):
        cell.text = header
    rows = [
        ['General', '0.8928', '0.8649', '0.8914', '0.8822', '0.8617'],
        ['Report', '0.9040', '0.8604', '0.8991', '0.8660', '0.8592'],
    ]
    for row in rows:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            cell.text = value
    for row in table.rows:
        for cell in row.cells:
            for p in cell.paragraphs:
                set_para_spacing(p, line_15=False)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in p.runs:
                    run.font.name = 'Arial'
                    run.font.size = Pt(10.5)


def build_dissertation_docx(md_path: Path, out_path: Path):
    doc = Document()
    set_document_layout(doc)

    # Title page
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(110)
    r = p.add_run('Design and Development of a Hybrid Chatbot for Medical Haematology using Sequential Models for Intent Classification and Retrieval-Based Response Generation')
    r.bold = True
    r.font.name = 'Arial'
    r.font.size = Pt(20)
    for line in [
        '',
        'Final Dissertation Submission Version',
        'Programme: [Insert Degree Title]',
        'Student Name: [Insert Name]',
        'Student ID: [Insert ID]',
        'Supervisor: [Insert Supervisor Name]',
        'Submission Year: 2026',
    ]:
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        para.add_run(line).font.name = 'Arial'
        set_para_spacing(para, line_15=False)
    doc.add_page_break()

    md_text = md_path.read_text(encoding='utf-8')
    abstract = md_text.split('## Abstract', 1)[1].split('## Chapter 1', 1)[0].strip().split('\n\n')
    add_body_paragraph(doc, 'Abstract', style='Heading 1', single_space=True)
    for para in abstract:
        add_body_paragraph(doc, para.strip(), single_space=True)

    doc.add_page_break()
    add_body_paragraph(doc, 'Acknowledgements', style='Heading 1', single_space=True)
    add_body_paragraph(doc, 'Support provided by the project supervisor, testers, and reviewers was acknowledged. Their comments were incorporated into the final system and documentation.', single_space=True)

    doc.add_page_break()
    add_body_paragraph(doc, 'Table of Contents', style='Heading 1', single_space=True)
    toc = doc.add_paragraph()
    add_toc(toc)
    set_para_spacing(toc, line_15=False)

    main_section = doc.add_section(WD_SECTION.NEW_PAGE)
    main_section.left_margin = Cm(3)
    main_section.right_margin = Cm(2)
    main_section.top_margin = Cm(2.5)
    main_section.bottom_margin = Cm(2.5)
    restart_page_numbering(main_section)
    footer = main_section.footer
    footer.is_linked_to_previous = False
    add_page_field(footer.paragraphs[0])

    main = md_text.split('## Chapter 1', 1)[1].split('## References', 1)[0]
    pending = []

    def flush_pending():
        nonlocal pending
        if pending:
            add_body_paragraph(doc, ' '.join(pending).strip())
            pending = []

    for raw_line in ('## Chapter 1' + main).splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            flush_pending()
            continue
        if line.startswith('## '):
            flush_pending()
            title = line[3:]
            add_body_paragraph(doc, title, style='Heading 1')
            if title == 'Chapter 3: Design and Development':
                add_figure(doc, FIGURE_MAP['Figure 3.1'], 'Figure 3.1 Hybrid haematology chatbot architecture', Inches(6.4))
                add_figure(doc, FIGURE_MAP['Figure 3.2'], 'Figure 3.2 MLOps lifecycle for the implemented chatbot', Inches(6.2))
                add_figure(doc, FIGURE_MAP['Figure 3.3'], 'Figure 3.3 Implemented chat user interface', Inches(6.2))
            if title == 'Chapter 4: Result Analysis and Evaluation':
                add_metrics_table(doc)
                add_figure(doc, FIGURE_MAP['Figure 4.1'], 'Figure 4.1 Test-set metric comparison across model profiles', Inches(6.1))
                add_figure(doc, FIGURE_MAP['Figure 4.2'], 'Figure 4.2 Training and validation behaviour of the two BiLSTM models', Inches(6.3))
                add_figure(doc, FIGURE_MAP['Figure 4.3'], 'Figure 4.3 Dataset footprint and fixed split allocation', Inches(6.2))
                add_figure(doc, FIGURE_MAP['Figure 4.4'], 'Figure 4.4 Automated test execution summary', Inches(5.8))
                add_figure(doc, FIGURE_MAP['Figure 4.5'], 'Figure 4.5 Admin overview dashboard for operational monitoring', Inches(6.3))
                add_figure(doc, FIGURE_MAP['Figure 4.6'], 'Figure 4.6 Inference trace for the phrase "What is aPTT?"', Inches(6.3))
            continue
        if line.startswith('### '):
            flush_pending()
            add_body_paragraph(doc, line[4:], style='Heading 2')
            continue
        if re.match(r'^\d+\.\s', line):
            flush_pending()
            add_body_paragraph(doc, re.sub(r'^\d+\.\s', '', line), style='List Number')
            continue
        if line.startswith('- '):
            flush_pending()
            add_body_paragraph(doc, line[2:], style='List Bullet')
            continue
        if line.startswith('[Insert '):
            flush_pending()
            continue
        pending.append(line)
    flush_pending()

    references = md_text.split('## References', 1)[1].split('## Appendices Guidance', 1)[0].strip().split('\n\n')
    appendices = md_text.split('## Appendices Guidance (for the final Word version)', 1)[1].strip().splitlines()
    add_body_paragraph(doc, 'References', style='Heading 1')
    for ref in references:
        p = doc.add_paragraph()
        p.add_run(clean_inline_markdown(ref.strip()))
        p.paragraph_format.left_indent = Cm(0.5)
        p.paragraph_format.first_line_indent = Cm(-0.5)
        set_para_spacing(p)
    add_body_paragraph(doc, 'Appendices', style='Heading 1')
    for line in appendices:
        if not line.strip():
            continue
        if line.startswith('- '):
            add_body_paragraph(doc, line[2:], style='List Bullet')
        else:
            add_body_paragraph(doc, clean_inline_markdown(line))
    doc.save(out_path)
    print(f'Saved {out_path}')


def build_simple_docx(md_path: Path, out_path: Path, title: str, image_sequence: list[tuple[str, Path]] | None = None):
    doc = Document()
    set_document_layout(doc)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(title)
    run.bold = True
    run.font.name = 'Arial'
    run.font.size = Pt(18)
    set_para_spacing(p)

    lines = md_path.read_text(encoding='utf-8').splitlines()
    image_index = 0
    for raw in lines:
        line = raw.rstrip()
        if not line.strip():
            continue
        if line.startswith('# '):
            continue
        if line.startswith('## '):
            add_body_paragraph(doc, line[3:], style='Heading 1')
        elif line.startswith('### '):
            add_body_paragraph(doc, line[4:], style='Heading 2')
        elif line.startswith('- [ ] '):
            add_body_paragraph(doc, '? ' + line[6:], style='List Bullet')
        elif line.startswith('- '):
            add_body_paragraph(doc, line[2:], style='List Bullet')
        elif re.match(r'^\d+\.\d+', line):
            add_body_paragraph(doc, line, style='List Bullet')
        elif re.match(r'^\d+\.\s', line):
            add_body_paragraph(doc, re.sub(r'^\d+\.\s', '', line), style='List Number')
        elif line.startswith('```'):
            continue
        else:
            add_body_paragraph(doc, clean_inline_markdown(line))
        if image_sequence and image_index < len(image_sequence):
            trigger, path = image_sequence[image_index]
            if trigger in line:
                add_figure(doc, path, trigger, Inches(6.0))
                image_index += 1
    doc.save(out_path)
    print(f'Saved {out_path}')


def add_textbox(slide, left, top, width, height, text, font_size=24, color=TEXT_DARK, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = PptPt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return box


def add_bullets(box, bullets, font_size=24, color=TEXT_DARK):
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = 'Arial'
        p.font.size = PptPt(font_size)
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = 8


def build_presentation():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def base_slide(title: str, subtitle: str | None = None, dark=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BG_DARK if dark else BG_LIGHT
        title_box = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.45), PptInches(12), PptInches(0.8))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = PptPt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255,255,255) if dark else ACCENT
        if subtitle:
            sub = slide.shapes.add_textbox(PptInches(0.72), PptInches(1.08), PptInches(11.8), PptInches(0.4))
            sp = sub.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.name = 'Arial'
            sp.font.size = PptPt(12)
            sp.font.color.rgb = RGBColor(210,230,226) if dark else TEXT_MUTED
        return slide

    slides_spec = []

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = BG_DARK
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_DARK; shape.line.color.rgb = BG_DARK
    add_textbox(slide, PptInches(0.9), PptInches(1.0), PptInches(11.3), PptInches(1.6), 'Hybrid Chatbot for Medical Haematology', 28, RGBColor(255,255,255), True)
    add_textbox(slide, PptInches(0.95), PptInches(2.05), PptInches(10.8), PptInches(1.2), 'Sequential models for intent classification and retrieval-based response generation', 18, RGBColor(214,232,228), False)
    add_textbox(slide, PptInches(0.95), PptInches(5.9), PptInches(6.0), PptInches(0.5), 'Final Year Dissertation Presentation', 16, RGBColor(214,232,228), False)
    add_textbox(slide, PptInches(0.95), PptInches(6.35), PptInches(7.2), PptInches(0.5), 'Student: [Insert Name]   Supervisor: [Insert Name]   Year: 2026', 12, RGBColor(214,232,228), False)
    slide.shapes.add_picture(str(FIG / 'figure_3_1_hybrid_architecture.png'), PptInches(7.8), PptInches(1.35), width=PptInches(4.7), height=PptInches(4.8))
    slides_spec.append({'title':'Hybrid Chatbot for Medical Haematology','subtitle':'Sequential models for intent classification and retrieval-based response generation','bullets':['Final Year Dissertation Presentation'], 'image': FIG / 'figure_3_1_hybrid_architecture.png'})

    # Slide 2
    slide = base_slide('Research Context and Project Aim')
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.5), PptInches(6.0), PptInches(4.8))
    add_bullets(box, [
        'General-purpose chatbots do not naturally satisfy laboratory requirements for scope control, auditability, and safe fallback behaviour.',
        'The project therefore targeted a bounded haematology laboratory assistant rather than a diagnostic system.',
        'The aim was to combine sequential intent classification with retrieval-based response generation and practical MLOps controls.'
    ], 20)
    slide.shapes.add_picture(str(FIG / 'figure_3_3_chat_ui.png'), PptInches(7.35), PptInches(1.55), width=PptInches(5.2))
    slides_spec.append({'title':'Research Context and Project Aim','bullets':['General-purpose chatbots do not naturally satisfy laboratory requirements for scope control, auditability, and safe fallback behaviour.','The project therefore targeted a bounded haematology laboratory assistant rather than a diagnostic system.','The aim was to combine sequential intent classification with retrieval-based response generation and practical MLOps controls.'], 'image': FIG / 'figure_3_3_chat_ui.png'})

    # Slide 3
    slide = base_slide('System Architecture and Backend Layers')
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(4.5), PptInches(5.0))
    add_bullets(box, [
        'FastAPI served the chat API, admin endpoints, and model registry.',
        'The inference path combined text preprocessing, BiLSTM intent scoring, entity detection, TF-IDF retrieval, and safety rules.',
        'PostgreSQL logging and the admin console supported monitoring and retraining.'
    ], 19)
    slide.shapes.add_picture(str(FIG / 'figure_3_1_hybrid_architecture.png'), PptInches(5.55), PptInches(1.35), width=PptInches(6.8))
    slides_spec.append({'title':'System Architecture and Backend Layers','bullets':['FastAPI served the chat API, admin endpoints, and model registry.','The inference path combined text preprocessing, BiLSTM intent scoring, entity detection, TF-IDF retrieval, and safety rules.','PostgreSQL logging and the admin console supported monitoring and retraining.'], 'image': FIG / 'figure_3_1_hybrid_architecture.png'})

    # Slide 4
    slide = base_slide('Data Strategy and Two-Model Design')
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(6.0), PptInches(4.9))
    add_bullets(box, [
        'Master labelled dataset: 3,250 utterances across 27 intents.',
        'General model dataset: 2,670 utterances across 21 intents.',
        'Report model dataset: 2,636 utterances across 24 intents.',
        'The split reduced label competition between workflow questions and CBC report-support questions.'
    ], 20)
    slide.shapes.add_picture(str(FIG / 'figure_4_3_dataset_and_splits.png'), PptInches(7.3), PptInches(1.55), width=PptInches(5.1))
    slides_spec.append({'title':'Data Strategy and Two-Model Design','bullets':['Master labelled dataset: 3,250 utterances across 27 intents.','General model dataset: 2,670 utterances across 21 intents.','Report model dataset: 2,636 utterances across 24 intents.','The split reduced label competition between workflow questions and CBC report-support questions.'], 'image': FIG / 'figure_4_3_dataset_and_splits.png'})

    # Slide 5
    slide = base_slide('MLOps Lifecycle and Operational Monitoring')
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(5.4), PptInches(5.0))
    add_bullets(box, [
        'The implemented lifecycle covered labelled data ingestion, derived datasets, fixed splits, training, evaluation, deployment, monitoring, and retraining.',
        'The admin console exposed model metrics, live logs, version-relevant artefacts, review workflow, and inference traces.',
        'Reviewed production phrases could be exported and merged back into the labelled dataset.'
    ], 19)
    slide.shapes.add_picture(str(FIG / 'figure_3_2_mlops_lifecycle.png'), PptInches(6.65), PptInches(1.3), width=PptInches(6.0))
    slides_spec.append({'title':'MLOps Lifecycle and Operational Monitoring','bullets':['The implemented lifecycle covered labelled data ingestion, derived datasets, fixed splits, training, evaluation, deployment, monitoring, and retraining.','The admin console exposed model metrics, live logs, version-relevant artefacts, review workflow, and inference traces.','Reviewed production phrases could be exported and merged back into the labelled dataset.'], 'image': FIG / 'figure_3_2_mlops_lifecycle.png'})

    # Slide 6
    slide = base_slide('Evaluation Design and Testing Discipline')
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(5.8), PptInches(5.0))
    add_bullets(box, [
        'A fixed 70/15/15 train-validation-test policy was adopted to separate checkpoint selection from final reporting.',
        'Automated tests covered API behaviour, preprocessing, entity detection, routing, model advisory, predictor rules, and split generation.',
        'The final automated suite passed with 26 tests.'
    ], 19)
    slide.shapes.add_picture(str(FIG / 'figure_4_4_testing_summary.png'), PptInches(7.0), PptInches(1.55), width=PptInches(5.2))
    slides_spec.append({'title':'Evaluation Design and Testing Discipline','bullets':['A fixed 70/15/15 train-validation-test policy was adopted to separate checkpoint selection from final reporting.','Automated tests covered API behaviour, preprocessing, entity detection, routing, model advisory, predictor rules, and split generation.','The final automated suite passed with 26 tests.'], 'image': FIG / 'figure_4_4_testing_summary.png'})

    # Slide 7
    slide = base_slide('Model Performance')
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(5.5), PptInches(5.0))
    add_bullets(box, [
        'General model: accuracy 0.8928, macro F1 0.8649.',
        'Report model: accuracy 0.9040, macro F1 0.8604.',
        'The report profile achieved slightly higher accuracy, while both models maintained near-0.90 test performance within a bounded medical scope.'
    ], 20)
    slide.shapes.add_picture(str(FIG / 'figure_4_1_metric_comparison.png'), PptInches(6.95), PptInches(1.7), width=PptInches(5.1))
    slides_spec.append({'title':'Model Performance','bullets':['General model: accuracy 0.8928, macro F1 0.8649.','Report model: accuracy 0.9040, macro F1 0.8604.','The report profile achieved slightly higher accuracy, while both models maintained near-0.90 test performance within a bounded medical scope.'], 'image': FIG / 'figure_4_1_metric_comparison.png'})

    # Slide 8
    slide = base_slide('Training Behaviour and Interpretability')
    box = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.45), PptInches(5.3), PptInches(5.0))
    add_bullets(box, [
        'Both models showed sustained loss reduction and validation improvement across the eight-epoch training schedule.',
        'The remaining gap between validation and test results indicated moderate but acceptable generalisation loss.',
        'The inference trace made tokenisation, top intent scores, entity routing, retrieval candidates, and final response selection visible.'
    ], 19)
    slide.shapes.add_picture(str(FIG / 'figure_4_2_training_curves.png'), PptInches(6.55), PptInches(1.2), width=PptInches(6.0))
    slides_spec.append({'title':'Training Behaviour and Interpretability','bullets':['Both models showed sustained loss reduction and validation improvement across the eight-epoch training schedule.','The remaining gap between validation and test results indicated moderate but acceptable generalisation loss.','The inference trace made tokenisation, top intent scores, entity routing, retrieval candidates, and final response selection visible.'], 'image': FIG / 'figure_4_2_training_curves.png'})

    # Slide 9
    slide = base_slide('Strengths, Limitations, and Risk Controls')
    box = slide.shapes.add_textbox(PptInches(0.75), PptInches(1.4), PptInches(12), PptInches(4.9))
    add_bullets(box, [
        'Strengths: bounded medical scope, controlled retrieval, reproducible splits, live monitoring, review-driven retraining.',
        'Limitations: moderate dataset size, class imbalance, English-only support, rule-based entity refinement, no patient-specific diagnosis or treatment.',
        'Risk control: unsafe prompts, out-of-scope prompts, and incomplete queries were handled through dedicated intents and guardrail responses.'
    ], 21)
    slides_spec.append({'title':'Strengths, Limitations, and Risk Controls','bullets':['Strengths: bounded medical scope, controlled retrieval, reproducible splits, live monitoring, review-driven retraining.','Limitations: moderate dataset size, class imbalance, English-only support, rule-based entity refinement, no patient-specific diagnosis or treatment.','Risk control: unsafe prompts, out-of-scope prompts, and incomplete queries were handled through dedicated intents and guardrail responses.'], 'image': None})

    # Slide 10
    slide = base_slide('Conclusion and Further Work', dark=True)
    box = slide.shapes.add_textbox(PptInches(0.9), PptInches(1.55), PptInches(6.1), PptInches(4.9))
    add_bullets(box, [
        'The project delivered a substantial artefact that combined sequential NLP, retrieval-based response generation, and practical MLOps controls.',
        'Further work should extend real-user data collection, multilingual support, richer report-term coverage, and formal usability evaluation.',
        'The implemented platform is suitable as a bounded laboratory assistant and as a demonstrable final-year machine learning system.'
    ], 20, RGBColor(236,245,243))
    slide.shapes.add_picture(str(FIG / 'figure_4_5_admin_overview.png'), PptInches(7.15), PptInches(1.35), width=PptInches(5.0))
    slides_spec.append({'title':'Conclusion and Further Work','subtitle':'The project delivered a substantial artefact that combined sequential NLP, retrieval-based response generation, and practical MLOps controls.','bullets':['Further work should extend real-user data collection, multilingual support, richer report-term coverage, and formal usability evaluation.'], 'image': FIG / 'figure_4_5_admin_overview.png', 'dark':True})

    prs.save(PPTX_OUT)
    print(f'Saved {PPTX_OUT}')

    # Simple preview PNGs for QA
    font_title = ImageFont.truetype('C:/Windows/Fonts/arialbd.ttf', 54)
    font_sub = ImageFont.truetype('C:/Windows/Fonts/arial.ttf', 24)
    font_body = ImageFont.truetype('C:/Windows/Fonts/arial.ttf', 30)
    font_small = ImageFont.truetype('C:/Windows/Fonts/arial.ttf', 18)
    for idx, spec in enumerate(slides_spec, start=1):
        dark = spec.get('dark', False)
        bg = (14,49,47) if dark else (247,250,249)
        fg = (255,255,255) if dark else (13,84,78)
        body = (235,245,243) if dark else (24,39,40)
        muted = (214,232,228) if dark else (81,98,100)
        img = Image.new('RGB', (1920,1080), bg)
        draw = ImageDraw.Draw(img)
        draw.text((105,70), spec['title'], fill=fg, font=font_title)
        if spec.get('subtitle'):
            draw.text((108,150), spec['subtitle'], fill=muted, font=font_sub)
        y = 250 if not spec.get('subtitle') else 200
        for bullet in spec.get('bullets', []):
            bullet_lines = wrap(bullet, width=42)
            draw.text((125, y), u'• ' + bullet_lines[0], fill=body, font=font_body)
            y += 44
            for line in bullet_lines[1:]:
                draw.text((165, y), line, fill=body, font=font_body)
                y += 40
            y += 20
        if spec.get('image'):
            source = Image.open(spec['image']).convert('RGB')
            source.thumbnail((860, 620))
            panel = Image.new('RGB', (900, 640), (255,255,255) if not dark else (28,69,66))
            px = (panel.width - source.width)//2
            py = (panel.height - source.height)//2
            panel.paste(source, (px, py))
            img.paste(panel, (960, 240))
        draw.text((105,1020), f'Slide {idx}', fill=muted, font=font_small)
        out = PREVIEW_DIR / f'slide_{idx:02d}.png'
        img.save(out)
    print(f'Preview images written to {PREVIEW_DIR}')


def main():
    build_dissertation_docx(DISSERTATION_MD, DISSERTATION_DOCX)
    build_simple_docx(CHECKLIST_MD, CHECKLIST_DOCX, 'Submission Checklist for Hybrid Haematology Chatbot Dissertation')
    build_simple_docx(MANUAL_MD, MANUAL_DOCX, 'User Manual for the Hybrid Haematology Chatbot', image_sequence=[
        ('## 7. Using the Chatbot', FIG / 'figure_3_3_chat_ui.png'),
        ('### 8.1 Overview', FIG / 'figure_4_5_admin_overview.png'),
        ('### 8.2 Inference Trace', FIG / 'figure_4_7_admin_trace.png'),
    ])
    build_presentation()

if __name__ == '__main__':
    main()
